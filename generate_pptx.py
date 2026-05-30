from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()

    def set_slide_title(slide, text):
        title = slide.shapes.title
        title.text = text
        for paragraph in title.text_frame.paragraphs:
            paragraph.font.bold = True
            paragraph.font.size = Pt(36)

    def add_bullet_points(slide, points):
        body_shape = slide.shapes.placeholders[1]
        tf = body_shape.text_frame
        for point in points:
            p = tf.add_paragraph()
            p.text = point
            p.level = 0
            p.font.size = Pt(20)

    # --- Slide 1: Title ---
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.shapes.placeholders[1]
    title.text = "StageLine: Ολοκληρωμένο Σύστημα Κρατήσεων Θεάτρου"
    subtitle.text = "Mobile App (React Native), REST API (Node.js) & MariaDB\n\nΕργασία στο μάθημα: Κατανεμημένα Συστήματα"

    # --- Slide 2: Εισαγωγή & Στόχος ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_title(slide, "Εισαγωγή & Στόχος")
    add_bullet_points(slide, [
        "Ανάπτυξη cross-platform εφαρμογής για θεατρικές κρατήσεις.",
        "Στόχος: Εύκολη εύρεση παραστάσεων και ασφαλής κράτηση θέσεων.",
        "Έμφαση στην ταυτοποίηση χρηστών και τη συνέπεια δεδομένων.",
        "Υποστήριξη real-time διαθεσιμότητας θέσεων."
    ])

    # --- Slide 3: Τεχνολογικό Stack ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_title(slide, "Τεχνολογικό Stack")
    add_bullet_points(slide, [
        "Frontend: React Native & Expo (NativeWind/Tailwind CSS).",
        "Backend: Node.js & Express (RESTful API).",
        "Database: MariaDB / MySQL.",
        "State Management: Zustand & React Query.",
        "Ασφάλεια: Bcrypt.js & JSON Web Tokens (JWT)."
    ])

    # --- Slide 4: Αρχιτεκτονική Συστήματος ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_title(slide, "Αρχιτεκτονική Συστήματος")
    add_bullet_points(slide, [
        "Client-Server Architecture.",
        "Mobile Client: React Native App (Expo).",
        "RESTful API: Middleware για Authentication & Validation.",
        "Data Layer: MariaDB για μόνιμη αποθήκευση.",
        "Fallback Mechanism: In-memory store αν η DB είναι offline."
    ])

    # --- Slide 5: Λειτουργίες Frontend ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_title(slide, "Λειτουργίες Frontend (Χρήστης)")
    add_bullet_points(slide, [
        "User Registration & Secure Login.",
        "Dynamic Search: Αναζήτηση ανά τίτλο, θέατρο ή τοποθεσία.",
        "Studio Network: Φιλτράρισμα παραστάσεων ανά σκηνή.",
        "User Profile: Διαχείριση στοιχείων και ιστορικό κρατήσεων."
    ])

    # --- Slide 6: Διαδικασία Κράτησης & Τροποποίησης ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_title(slide, "Κρατήσεις & Τροποποίηση")
    add_bullet_points(slide, [
        "Interactive Seat Map: Επιλογή θέσεων με οπτική ανατροφοδότηση.",
        "Real-time Booking: Επιβεβαίωση κράτησης μέσω API.",
        "Modify Slot: Δυνατότητα αλλαγής θέσης/ώρας σε μελλοντικές κρατήσεις.",
        "Cancellation: Πλήρης διαγραφή και απελευθέρωση θέσης."
    ])

    # --- Slide 7: Ασφάλεια & Ταυτοποίηση ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_title(slide, "Ασφάλεια & Ταυτοποίηση")
    add_bullet_points(slide, [
        "JWT Strategy: Access Tokens (15m) & Refresh Tokens (30d).",
        "Secure Storage: Χρήση expo-secure-store στη συσκευή.",
        "Axios Interceptors: Αυτόματη ανανέωση session (silent refresh).",
        "Password Hashing: Προστασία κωδικών με Salted Bcrypt."
    ])

    # --- Slide 8: Backend & API Endpoints ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_title(slide, "Backend & API Endpoints")
    add_bullet_points(slide, [
        "Auth: /register, /login, /refresh, /me.",
        "Content: /theatres, /shows, /showtimes, /seats.",
        "Reservations: POST (Create), PATCH (Update), DELETE (Cancel).",
        "User History: /user/reservations."
    ])

    # --- Slide 9: Σχεδιασμός Βάσης Δεδομένων ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_title(slide, "Βάση Δεδομένων (MariaDB)")
    add_bullet_points(slide, [
        "Users Table: Αποθήκευση hashed credentials.",
        "Theatres & Shows: Σχέση One-to-Many.",
        "Showtimes: Διαχείριση τιμών και ωραρίων.",
        "Reservations: Σύνδεση User <-> Showtime με unique seat constraint.",
        "Auto-schema generation κατά την εκκίνηση."
    ])

    # --- Slide 10: UI/UX Σχεδιασμός ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_title(slide, "UI/UX Σχεδιασμός")
    add_bullet_points(slide, [
        "Dark Mode Aesthetics: Χρήση Midnight Indigo και Electric Cyan.",
        "Atmospheric Aurora Background: Δυναμικά οπτικά εφέ.",
        "Responsive Grid: Προσαρμογή σε διαφορετικά μεγέθη οθονών.",
        "Feedback: Alerts και Indicators για κάθε ενέργεια του χρήστη."
    ])

    # --- Slide 11: Καινοτομίες & Fallback ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_title(slide, "Καινοτομίες & Fallback")
    add_bullet_points(slide, [
        "Hybrid Backend: Αυτόματη εναλλαγή σε Memory Store αν η DB αποτύχει.",
        "Smart Search: Αναζήτηση σε πολλαπλά πεδία (Title/Theatre/Location).",
        "Modular Code: Διαχωρισμός Routes, Controllers και Store logic.",
        "Clean Architecture: Εύκολη επέκταση και συντήρηση."
    ])

    # --- Slide 12: Επίλογος ---
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.shapes.placeholders[1]
    title.text = "Ευχαριστούμε για την προσοχή σας!"
    subtitle.text = "Ερωτήσεις;"

    prs.save('StageLine_Presentation.pptx')

if __name__ == "__main__":
    create_presentation()
